from flask import Flask, request, jsonify, Response
from flask_cors import CORS  # Added for CORS support
import os
import requests
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import io
from weasyprint import HTML, CSS
import markdown
import pypandoc
from docx.shared import Pt
from docx.oxml.ns import qn
import base64

app = Flask(__name__)

# Enable CORS only for the relevant routes
CORS(app, resources={r"/search-sober-living": {"origins": "*"}, r"/search-sober-living/get-details": {"origins": "*"}})

@app.route('/format-html', methods=['POST'])
def format_html():
    # Get HTML from form-data
    html_code = request.form.get("html", "")

    if not html_code:
        return jsonify({"error": "No HTML provided"}), 400

    # Parse and prettify HTML
    soup = BeautifulSoup(html_code, "html.parser")
    formatted_html = soup.prettify()

    return jsonify({"formatted_html": formatted_html})
    
@app.route("/image-to-base64", methods=["POST"])
def image_to_base64():
    try:
        url = request.form.get("url")
        if not url:
            return jsonify(error="No URL provided"), 400

        response = requests.get(url)
        response.raise_for_status()

        image_base64 = base64.b64encode(response.content).decode("utf-8")

        return jsonify(base64=image_base64), 200

    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/convert', methods=['POST'])
def convert():
    if 'url' not in request.form:
        return jsonify(error="No URL provided"), 400

    url = request.form['url']

    if url.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
        text, error = extract_text_from_image(url)
    else:
        text, error = extract_text_from_url(url)

    if error:
        return jsonify(error=error), 500
    else:
        return jsonify(text=text), 200

def preprocess_markdown(md_content):
    # Ensure each line becomes a paragraph
    lines = md_content.splitlines()
    return "\n\n".join([line.strip() for line in lines if line.strip()])

@app.route("/convert-markdown/pdf", methods=["POST"])
def convert_markdown_pdf():
    try:
        markdown_content = request.form.get("markdown")
        if not markdown_content:
            return jsonify(error="No Markdown provided"), 400

        # Convert Markdown → HTML with proper line breaks and tables
        html_content = markdown.markdown(
            markdown_content,
            extensions=['tables', 'extra', 'sane_lists', 'nl2br']
        )

        # Add some CSS for proper spacing and font
        css = """
        body { font-family: Arial, sans-serif; line-height: 1.5; }
        p { margin: 0 0 10px; }
        ul, ol { margin: 0 0 10px 20px; }
        li { margin-bottom: 5px; }
        table { border-collapse: collapse; margin-bottom: 10px; }
        th, td { border: 1px solid #333; padding: 5px; }
        """

        pdf_bytes = HTML(string=html_content).write_pdf(stylesheets=[CSS(string=css)])

        return Response(
            pdf_bytes,
            mimetype="application/pdf",
            headers={"Content-Disposition": "attachment; filename=output.pdf"}
        )

    except Exception as e:
        return jsonify(error=str(e)), 500
        
@app.route("/convert-markdown/docx", methods=["POST"])
def convert_markdown_docx():
    try:
        markdown_content = request.form.get("markdown")
        if not markdown_content:
            return jsonify(error="No Markdown provided"), 400

        md_for_docx = preprocess_markdown(markdown_content)

        temp_docx = "temp.docx"
        pypandoc.convert_text(
            md_for_docx,
            "docx",
            format="md",
            outputfile=temp_docx,
            extra_args=["--standalone", "--from=markdown+pipe_tables"]
        )

        doc = Document(temp_docx)
        for para in doc.paragraphs:
            for run in para.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(12)

        docx_io = io.BytesIO()
        doc.save(docx_io)
        docx_bytes = docx_io.getvalue()

        return Response(
            docx_bytes,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=output.docx"}
        )

    except Exception as e:
        return jsonify(error=str(e)), 500
        
@app.route('/url-text-extract', methods=['POST'])
def url_text_extract():
    if 'url' not in request.form:
        return jsonify(error="No URL provided"), 400

    url = request.form['url']

    text, error = extract_text_from_webpage(url)

    if error:
        return jsonify(error=error), 500
    else:
        return jsonify(text=text), 200

@app.route('/search-sober-living', methods=['GET'])
def search_sober_living():
    query = request.args.get('query')
    api_key = request.args.get('api_key')

    if not query:
        return jsonify(error="No query provided"), 400
    if not api_key:
        return jsonify(error="No API Key provided"), 400

    url = f'https://maps.googleapis.com/maps/api/place/textsearch/json?query={query}&key={api_key}'
    try:
        response = requests.get(url)
        data = response.json()

        if response.status_code == 200 and data.get('status') == 'OK':
            return jsonify(data), 200
        else:
            return jsonify(error=f"Google API Error: {data.get('status', 'Unknown Error')}"), 500
    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/search-sober-living/get-details', methods=['GET'])
def get_place_details():
    place_id = request.args.get('place_id')
    api_key = request.args.get('api_key')

    if not place_id:
        return jsonify(error="No place_id provided"), 400
    if not api_key:
        return jsonify(error="No API Key provided"), 400

    url = f'https://maps.googleapis.com/maps/api/place/details/json?place_id={place_id}&key={api_key}'
    try:
        response = requests.get(url)
        data = response.json()

        if response.status_code == 200 and data.get('status') == 'OK':
            return jsonify(result=data['result']), 200
        else:
            return jsonify(error=f"Google API Error: {data.get('status', 'Unknown Error')}"), 500
    except Exception as e:
        return jsonify(error=str(e)), 500

@app.route('/search-sober-living/next-page', methods=['GET'])
def search_sober_living_next_page():
    next_page_token = request.args.get('next_page_token')
    api_key = request.args.get('api_key')

    if not next_page_token:
        return jsonify(error="No next_page_token provided"), 400
    if not api_key:
        return jsonify(error="No API Key provided"), 400

    url = f'https://maps.googleapis.com/maps/api/place/textsearch/json?pagetoken={next_page_token}&key={api_key}'
    try:
        response = requests.get(url)
        data = response.json()

        if response.status_code == 200 and data.get('status') == 'OK':
            return jsonify(data), 200
        else:
            return jsonify(error=f"Google API Error: {data.get('status', 'Unknown Error')}"), 500
    except Exception as e:
        return jsonify(error=str(e)), 500

def extract_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        _, extension = os.path.splitext(url)
        extension = extension[1:].lower()

        if extension in ['pdf', 'docx', 'txt']:
            if extension == 'pdf':
                reader = PdfReader(io.BytesIO(response.content))
                text = " ".join(page.extract_text() for page in reader.pages)
            elif extension == 'docx':
                doc = Document(io.BytesIO(response.content))
                text = " ".join(paragraph.text for paragraph in doc.paragraphs)
            elif extension == 'txt':
                text = response.text
        elif extension == 'pptx':
            prs = Presentation(io.BytesIO(response.content))
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
        else:
            return None, "Unsupported filetype"

        return text, None
    except Exception as e:
        return None, str(e)

def extract_text_from_webpage(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, 'html.parser')
        text = '\n'.join(p.get_text() for p in soup.find_all('p'))

        return text, None
    except Exception as e:
        return None, str(e)

def extract_text_from_image(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        image = Image.open(io.BytesIO(response.content))
        text = pytesseract.image_to_string(image)

        return text, None
    except Exception as e:
        return None, str(e)

if __name__ == "__main__":
    app.run(debug=True)
